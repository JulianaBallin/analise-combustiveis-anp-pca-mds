from __future__ import annotations

import re
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_JUSTIFY
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import BaseDocTemplate, Frame, Image, PageTemplate, Paragraph, Spacer

from src.data_preparation import (
    FEATURES_NUMERICAS,
    PRECO_URL,
    VENDAS_URL,
    padronizar_features,
    preparar_dados,
)
from src.mds_analysis import aplicar_mds
from src.pca_analysis import aplicar_pca


ROOT = Path(__file__).resolve().parent
OUT_FIG = ROOT / "outputs" / "figures"
OUT_TAB = ROOT / "outputs" / "tables"
DATA_PROCESSED = ROOT / "data" / "processed"
REPORT_DIR = ROOT / "docs" / "report"
SLIDES_DIR = ROOT / "docs" / "slides"
NOTEBOOK_DIR = ROOT / "notebooks"
LOGO_PNG = ROOT / "docs" / "diagrams" / "logo.png"
PERIODO_INICIO = 2021
PERIODO_FIM = 2025
MAX_REGISTROS_MDS = 800


def main() -> None:
    configurar_estilo()
    criar_diretorios()

    dados = preparar_dados(periodo_inicio=PERIODO_INICIO, periodo_fim=PERIODO_FIM)
    dados_padronizados, _ = padronizar_features(dados)

    pca_df, cargas, variancia, _ = aplicar_pca(dados_padronizados)
    pca_plot = dados.join(pca_df)

    mds_df, mds = aplicar_mds(dados_padronizados, max_registros=MAX_REGISTROS_MDS)
    mds_plot = dados.loc[mds_df.index].join(mds_df)

    outliers = identificar_outliers(pca_plot)
    resumo_uf = resumo_por_uf(dados)

    salvar_tabelas(dados, variancia, cargas, outliers, resumo_uf, mds)
    gerar_figuras(pca_plot, mds_plot, variancia, cargas, outliers)

    contexto = montar_contexto(dados, variancia, cargas, outliers, resumo_uf, mds)
    gerar_relatorio_md(contexto)
    gerar_relatorio_pdf()
    gerar_slides(contexto)
    gerar_roteiro_slides_md(contexto)
    gerar_notebook()


def criar_diretorios() -> None:
    for path in [OUT_FIG, OUT_TAB, DATA_PROCESSED, REPORT_DIR, SLIDES_DIR, NOTEBOOK_DIR]:
        path.mkdir(parents=True, exist_ok=True)


def configurar_estilo() -> None:
    sns.set_theme(style="whitegrid", context="notebook")
    plt.rcParams["figure.dpi"] = 130
    plt.rcParams["savefig.bbox"] = "tight"


def identificar_outliers(pca_plot: pd.DataFrame) -> pd.DataFrame:
    dados = pca_plot.copy()
    dados["distancia_pca"] = (dados["PC1"] ** 2 + dados["PC2"] ** 2) ** 0.5
    colunas = [
        "mes_ano",
        "uf",
        "regiao",
        "preco_medio_gasolina_c",
        "volume_gasolina_c_m3",
        "volume_etanol_hidratado_m3",
        "participacao_etanol",
        "variacao_preco_gasolina_c",
        "variacao_volume_gasolina_c",
        "PC1",
        "PC2",
        "distancia_pca",
    ]
    return dados.sort_values("distancia_pca", ascending=False)[colunas].head(12)


def resumo_por_uf(dados: pd.DataFrame) -> pd.DataFrame:
    return (
        dados.groupby(["regiao", "uf"], as_index=False)
        .agg(
            preco_medio_gasolina_c=("preco_medio_gasolina_c", "mean"),
            volume_medio_gasolina_c_m3=("volume_gasolina_c_m3", "mean"),
            volume_medio_etanol_hidratado_m3=("volume_etanol_hidratado_m3", "mean"),
            participacao_media_etanol=("participacao_etanol", "mean"),
            variacao_media_preco=("variacao_preco_gasolina_c", "mean"),
            variacao_media_volume=("variacao_volume_gasolina_c", "mean"),
        )
        .sort_values("participacao_media_etanol", ascending=False)
    )


def salvar_tabelas(
    dados: pd.DataFrame,
    variancia: pd.DataFrame,
    cargas: pd.DataFrame,
    outliers: pd.DataFrame,
    resumo_uf: pd.DataFrame,
    mds,
) -> None:
    dados.to_csv(DATA_PROCESSED / "dataset_anp_pca_mds_2021_2025.csv", index=False)
    variancia.to_csv(OUT_TAB / "pca_variancia_explicada.csv", index=False)
    cargas.to_csv(OUT_TAB / "pca_cargas_componentes.csv")
    outliers.to_csv(OUT_TAB / "outliers_pca.csv", index=False)
    resumo_uf.to_csv(OUT_TAB / "resumo_por_uf.csv", index=False)
    pd.DataFrame({"metrica": ["stress_mds"], "valor": [float(mds.stress_)]}).to_csv(
        OUT_TAB / "mds_stress.csv", index=False
    )


def gerar_figuras(
    pca_plot: pd.DataFrame,
    mds_plot: pd.DataFrame,
    variancia: pd.DataFrame,
    cargas: pd.DataFrame,
    outliers: pd.DataFrame,
) -> None:
    plt.figure(figsize=(10, 7))
    sns.scatterplot(data=pca_plot, x="PC1", y="PC2", hue="regiao", style="regiao", s=48, alpha=0.8)
    plt.title("PCA 2D: estados e meses por região")
    plt.xlabel(f"PC1 ({variancia.loc[0, 'variancia_explicada']:.1%} da variância)")
    plt.ylabel(f"PC2 ({variancia.loc[1, 'variancia_explicada']:.1%} da variância)")
    plt.legend(title="Região", bbox_to_anchor=(1.02, 1), loc="upper left")
    plt.savefig(OUT_FIG / "pca_2d_regiao.png")
    plt.close()

    plt.figure(figsize=(10, 7))
    sns.scatterplot(data=mds_plot, x="MDS1", y="MDS2", hue="regiao", style="regiao", s=48, alpha=0.8)
    plt.title("MDS 2D: proximidade entre registros")
    plt.xlabel("Dimensão MDS 1")
    plt.ylabel("Dimensão MDS 2")
    plt.legend(title="Região", bbox_to_anchor=(1.02, 1), loc="upper left")
    plt.savefig(OUT_FIG / "mds_2d_regiao.png")
    plt.close()

    plt.figure(figsize=(8, 5))
    sns.barplot(data=variancia, x="componente", y="variancia_explicada", color="#287c8e")
    plt.ylim(0, 1)
    plt.title("Variância explicada no PCA")
    plt.ylabel("Proporção da variância")
    plt.xlabel("Componente")
    plt.savefig(OUT_FIG / "pca_variancia_explicada.png")
    plt.close()

    cargas_plot = cargas.reset_index(names="feature").melt(
        id_vars="feature", var_name="componente", value_name="carga"
    )
    plt.figure(figsize=(10, 6))
    sns.barplot(data=cargas_plot, y="feature", x="carga", hue="componente")
    plt.axvline(0, color="black", linewidth=0.8)
    plt.title("Cargas das variáveis nos dois primeiros componentes")
    plt.xlabel("Carga")
    plt.ylabel("Feature")
    plt.savefig(OUT_FIG / "pca_cargas_componentes.png")
    plt.close()

    plt.figure(figsize=(10, 6))
    plot = outliers.copy()
    plot["registro"] = plot["uf"] + " " + plot["mes_ano"].dt.strftime("%Y-%m")
    sns.barplot(data=plot, y="registro", x="distancia_pca", color="#c75d2c")
    plt.title("Registros mais afastados na projeção PCA")
    plt.xlabel("Distância ao centro da projeção")
    plt.ylabel("Registro")
    plt.savefig(OUT_FIG / "outliers_pca.png")
    plt.close()


def montar_contexto(
    dados: pd.DataFrame,
    variancia: pd.DataFrame,
    cargas: pd.DataFrame,
    outliers: pd.DataFrame,
    resumo_uf: pd.DataFrame,
    mds,
) -> dict:
    pc1_top = cargas["PC1"].abs().sort_values(ascending=False).head(4).index.tolist()
    pc2_top = cargas["PC2"].abs().sort_values(ascending=False).head(4).index.tolist()
    maior_etanol = resumo_uf.head(5)
    menor_etanol = resumo_uf.tail(5).sort_values("participacao_media_etanol")
    periodo = f"{dados['mes_ano'].min():%m/%Y} a {dados['mes_ano'].max():%m/%Y}"
    return {
        "n_registros": len(dados),
        "periodo": periodo,
        "ufs": dados["uf"].nunique(),
        "variancia_pc1": variancia.loc[0, "variancia_explicada"],
        "variancia_pc2": variancia.loc[1, "variancia_explicada"],
        "variancia_total": variancia.loc[1, "variancia_acumulada"],
        "pc1_top": pc1_top,
        "pc2_top": pc2_top,
        "outliers": outliers,
        "maior_etanol": maior_etanol,
        "menor_etanol": menor_etanol,
        "stress_mds": float(mds.stress_),
    }


def gerar_relatorio_md(ctx: dict) -> None:
    maior_etanol_txt = ", ".join(
        f"{row.uf} ({row.participacao_media_etanol:.1%})"
        for row in ctx["maior_etanol"].itertuples()
    )
    menor_etanol_txt = ", ".join(
        f"{row.uf} ({row.participacao_media_etanol:.1%})"
        for row in ctx["menor_etanol"].itertuples()
    )
    outliers_txt = "\n".join(
        f"- {row.uf} em {row.mes_ano:%m/%Y}: gasolina R$ {row.preco_medio_gasolina_c:.2f}, "
        f"participação do etanol {row.participacao_etanol:.1%}, variação de volume da gasolina {row.variacao_volume_gasolina_c:.1%}."
        for row in ctx["outliers"].head(6).itertuples()
    )
    texto = f"""# Relatório curto: PCA e MDS em dados da ANP

## 1. Problema investigado

O projeto investiga como as variações no preço médio da gasolina C se associam ao volume mensal vendido de gasolina C e à participação do etanol hidratado nos estados brasileiros. O cenário de negócio é o de uma rede de postos ou distribuidora que precisa decidir estoque, mix comercial, campanhas e metas regionais sem depender apenas do faturamento agregado.

## 2. Explicação teórica do PCA

PCA, ou Análise de Componentes Principais, é uma técnica de redução de dimensionalidade que transforma variáveis numéricas correlacionadas em componentes principais. Cada componente é uma combinação linear das variáveis originais. O primeiro componente concentra a maior variância possível dos dados; o segundo concentra a maior parte da variância restante, mantendo ortogonalidade em relação ao primeiro.

No projeto, o PCA foi usado para projetar oito features numéricas em duas dimensões, observar tendências, identificar registros afastados e medir quais variáveis mais contribuíram para os primeiros componentes.

## 3. Explicação teórica do MDS

MDS, ou Escalonamento Multidimensional, é uma técnica de visualização que parte das distâncias entre registros. Seu objetivo é posicionar pontos em poucas dimensões preservando, tanto quanto possível, as proximidades do espaço original. Registros parecidos ficam próximos; registros diferentes tendem a ficar afastados.

Ao contrário do PCA, o MDS não informa diretamente a contribuição de cada feature. Por isso, ele foi usado principalmente para comparar similaridades, agrupamentos e outliers com a leitura obtida no PCA.

## 4. Dataset utilizado

Foram usadas duas bases públicas da ANP:

- Série histórica mensal do levantamento de preços por estado, com preço médio de revenda da gasolina C e do etanol hidratado.
- Vendas de derivados de petróleo e etanol, com volumes mensais vendidos em metros cúbicos por UF e produto.

O recorte final contém {ctx["n_registros"]} registros, {ctx["ufs"]} UFs e período de {ctx["periodo"]}. Cada registro representa uma combinação UF-mês após cruzamento entre preços e vendas.

Fontes:

- {PRECO_URL}
- {VENDAS_URL}

## 5. Features selecionadas e justificativa

As features numéricas analisadas foram:

- `preco_medio_gasolina_c`: preço central da pergunta de negócio.
- `volume_gasolina_c_m3`: demanda mensal do combustível principal.
- `volume_etanol_hidratado_m3`: demanda do possível substituto.
- `variacao_preco_gasolina_c`: mede choques ou quedas mensais de preço.
- `variacao_volume_gasolina_c`: mede a mudança mensal da demanda de gasolina.
- `participacao_etanol`: indica o peso do etanol no total gasolina C + etanol.
- `razao_etanol_gasolina`: compara diretamente o volume de etanol com o de gasolina.
- `preco_relativo_etanol_gasolina`: compara o preço do etanol com o da gasolina.

A UF e a região foram mantidas como variáveis categóricas para colorir os gráficos e apoiar a interpretação visual.

## 6. Preparação e padronização dos dados

As bases foram carregadas nos links da ANP, filtramos gasolina C e etanol hidratado, agregamos por mês e UF e fizemos o **join só com `mes_ano` e a sigla `uf`**. Quando o nome completo do estado ou a grafia da região divergia entre preços e vendas, especialmente nos rótulos de **Centro-Oeste** (`Centro Oeste` x `Centro-Oeste`), a versão antiga do pipeline chegava a derrubar em torno de **240 registros** de DF, GO, MS e MT. Neste fluxo novo isso não acontece porque não precisamos de `uf_nome` idêntico byte a byte só para cruzar. `uf_nome` e `região` vêm de um `combine_first` depois do merge e ainda passam por `normalizar_regiao`. Na sequência entram participação do etanol, razão volume, preço relativo e variações percentuais mensais dentro de cada UF.

Como as features têm escalas muito diferentes, todas as variáveis numéricas foram padronizadas com `StandardScaler`. Essa etapa impede que volumes em m³ dominem indevidamente preços em reais ou indicadores percentuais.

## 7. PCA 2D, variância explicada e contribuição das variáveis

Os dois primeiros componentes explicaram {ctx["variancia_total"]:.1%} da variância total: PC1 explicou {ctx["variancia_pc1"]:.1%} e PC2 explicou {ctx["variancia_pc2"]:.1%}.

As variáveis que mais influenciaram PC1 foram: {", ".join(ctx["pc1_top"])}. As que mais influenciaram PC2 foram: {", ".join(ctx["pc2_top"])}.

Isso indica que a primeira dimensão separou principalmente estados e meses pelo porte e composição do consumo, enquanto a segunda dimensão destacou mudanças relativas de preço, volume e competitividade do etanol.

Figura: `outputs/figures/pca_variancia_explicada.png`.

Figura: `outputs/figures/pca_cargas_componentes.png`.

Figura principal: `outputs/figures/pca_2d_regiao.png`.

## 8. MDS 2D e interpretação das proximidades

O MDS foi aplicado sobre os dados padronizados, usando distância euclidiana e amostra de até {MAX_REGISTROS_MDS} registros para manter a visualização viável. O stress calculado foi {ctx["stress_mds"]:.2f}; quanto menor esse valor, melhor a preservação das distâncias na projeção.

No gráfico MDS, registros próximos representam UFs e meses com perfis semelhantes de preço, volume, participação do etanol, razão etanol/gasolina e variações mensais. A leitura visual mostrou proximidade entre registros de comportamento regional semelhante e maior afastamento de estados com peso muito alto de etanol ou volumes muito superiores ao restante do país.

Figura principal: `outputs/figures/mds_2d_regiao.png`.

## 9. Comparação entre PCA e MDS

O PCA foi mais útil para explicar quais variáveis estruturam a projeção, pois fornece variância explicada e cargas dos componentes. O MDS foi mais útil para visualizar similaridade entre registros sem exigir que a estrutura principal seja linear.

Os padrões gerais foram compatíveis: estados com maior participação do etanol e grande escala de vendas tenderam a se destacar; registros de menor escala ou com participação baixa do etanol ficaram em áreas mais próximas entre si. A diferença é que o PCA facilita a explicação matemática, enquanto o MDS facilita a leitura de vizinhança.

## 10. Padrões identificados

Antes da análise, esperava-se encontrar diferenças regionais, destaque de estados com maior consumo absoluto, separação parcial de mercados onde o etanol tem maior relevância e alguns meses atípicos ligados a variações bruscas de preço ou volume.

Os estados com maior participação média do etanol no recorte foram: {maior_etanol_txt}. Os estados com menor participação média foram: {menor_etanol_txt}.

Esses padrões fazem sentido no contexto do projeto, porque o etanol é mais competitivo e mais presente em alguns mercados, enquanto em outros a gasolina C domina a composição de consumo.

## 11. Outliers encontrados

Os principais registros afastados no PCA foram:

{outliers_txt}

Esses outliers são explicados principalmente por combinações de grande escala de volume, participação elevada do etanol, variações mensais fortes ou preço relativo do etanol muito diferente do padrão nacional.

Figura: `outputs/figures/outliers_pca.png`.

## 12. Riscos de interpretação

Existe risco de interpretar agrupamentos visuais como prova causal. PCA e MDS mostram associação e proximidade, mas não provam que o aumento de preço causou diretamente queda de volume. Também há risco de perder informação ao reduzir oito variáveis para duas dimensões, além de limitações das bases agregadas da ANP e ausência de variáveis externas como renda, frota, inflação, câmbio, clima e campanhas comerciais.

## 13. Uso para redução de features

O PCA pode ser usado como apoio para redução de features, especialmente se um modelo futuro precisar condensar variáveis correlacionadas de preço, volume e participação. Porém, como os dois primeiros componentes não preservam toda a variância, uma redução agressiva para apenas duas dimensões deve ser feita com cautela.

O MDS é mais adequado para visualização exploratória do que para redução de features em modelos preditivos, porque suas dimensões não têm interpretação direta por variável.

## 14. Conclusão

PCA e MDS ajudaram a compreender que o mercado analisado não se organiza apenas pelo preço da gasolina C. A escala de vendas, a presença relativa do etanol, a razão entre etanol e gasolina e as variações mensais também estruturam os padrões observados.

Para o cenário da rede de postos ou distribuidora, os métodos ajudam a identificar UFs e meses com comportamento semelhante, encontrar registros atípicos e apoiar decisões de estoque, mix de produtos e campanhas regionais. A análise deve ser vista como etapa exploratória robusta, não como prova causal definitiva.
"""
    (REPORT_DIR / "relatorio_pca_mds_anp.md").write_text(texto, encoding="utf-8")


def _cabecalho_rodape(canvas, doc) -> None:
    canvas.saveState()
    page_w, page_h = A4
    cor_linha = colors.HexColor("#287c8e")
    cor_texto = colors.HexColor("#546a70")
    if LOGO_PNG.exists():
        canvas.drawImage(
            str(LOGO_PNG),
            1.5 * cm,
            page_h - 2.6 * cm,
            width=9 * cm,
            height=1.7 * cm,
            preserveAspectRatio=True,
            mask="auto",
        )
    canvas.setStrokeColor(cor_linha)
    canvas.setLineWidth(0.5)
    canvas.line(1.5 * cm, page_h - 2.9 * cm, page_w - 1.5 * cm, page_h - 2.9 * cm)
    canvas.line(1.5 * cm, 1.9 * cm, page_w - 1.5 * cm, 1.9 * cm)
    canvas.setFont("Helvetica", 9)
    canvas.setFillColor(cor_texto)
    canvas.drawCentredString(page_w / 2, 1.1 * cm, f"Página {doc.page}")
    canvas.restoreState()


def gerar_relatorio_pdf() -> None:
    md = (REPORT_DIR / "relatorio_pca_mds_anp.md").read_text(encoding="utf-8")

    styles = getSampleStyleSheet()
    styles["BodyText"].alignment = TA_JUSTIFY
    styles["BodyText"].leading = 14
    styles["BodyText"].spaceAfter = 4

    story = []
    for bloco in md.split("\n\n"):
        bloco = bloco.strip()
        if not bloco:
            continue
        if bloco.startswith("# "):
            story.append(Paragraph(bloco[2:], styles["Title"]))
        elif bloco.startswith("## "):
            story.append(Paragraph(bloco[3:], styles["Heading2"]))
        elif re.match(r"^Figura(?: principal)?:", bloco):
            match = re.search(r"`(outputs/figures/[^`]+)`", bloco)
            if match:
                fig_path = ROOT / match.group(1)
                if fig_path.exists():
                    img = Image(str(fig_path), width=14 * cm, height=8.5 * cm)
                    story.append(Spacer(1, 4))
                    story.append(img)
        else:
            html = bloco.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            html = html.replace("\n", "<br/>")
            story.append(Paragraph(html, styles["BodyText"]))
        story.append(Spacer(1, 8))

    top = 3.5 * cm
    bottom = 2.8 * cm
    left = right = 2.0 * cm
    doc = BaseDocTemplate(
        str(REPORT_DIR / "relatorio_pca_mds_anp.pdf"),
        pagesize=A4,
        topMargin=top,
        bottomMargin=bottom,
        leftMargin=left,
        rightMargin=right,
    )
    frame = Frame(left, bottom, A4[0] - left - right, A4[1] - top - bottom, id="main")
    doc.addPageTemplates([PageTemplate(id="main", frames=frame, onPage=_cabecalho_rodape)])
    doc.build(story)


def add_slide(prs: Presentation, titulo: str, bullets: list[str], imagem: Path | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = titulo
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(5.1 if imagem else 8.4), Inches(4.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.level = 0
    if imagem:
        slide.shapes.add_picture(str(imagem), Inches(5.9), Inches(1.3), width=Inches(3.6))


def gerar_slides(ctx: dict) -> None:
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "PCA e MDS em dados da ANP"
    title.placeholders[1].text = "Gasolina C, etanol hidratado, preços e volumes por UF"

    add_slide(
        prs,
        "Problema investigado",
        [
            "Entender associações entre preço da gasolina C, volume vendido e participação do etanol.",
            "Apoiar decisões de estoque, mix comercial, campanhas e metas regionais.",
            "Tratar a análise como exploração visual e estatística, sem inferir causalidade direta.",
        ],
    )
    add_slide(
        prs,
        "Dataset e preparação",
        [
            f"Recorte: {ctx['periodo']}, {ctx['ufs']} UFs e {ctx['n_registros']} registros UF-mês.",
            "Fontes: série histórica mensal de preços e vendas mensais de derivados e etanol da ANP.",
            "Features padronizadas com StandardScaler para equilibrar escalas diferentes.",
        ],
    )
    add_slide(
        prs,
        "PCA",
        [
            f"PC1 explicou {ctx['variancia_pc1']:.1%} e PC2 explicou {ctx['variancia_pc2']:.1%}.",
            f"Variância acumulada dos dois primeiros componentes: {ctx['variancia_total']:.1%}.",
            f"Principais influências em PC1: {', '.join(ctx['pc1_top'])}.",
            f"Principais influências em PC2: {', '.join(ctx['pc2_top'])}.",
        ],
        OUT_FIG / "pca_2d_regiao.png",
    )
    add_slide(
        prs,
        "MDS",
        [
            "Representa proximidades entre registros usando distâncias no espaço padronizado.",
            f"Stress do MDS: {ctx['stress_mds']:.2f}.",
            "Registros próximos indicam comportamento semelhante de preço, volume e participação do etanol.",
        ],
        OUT_FIG / "mds_2d_regiao.png",
    )
    add_slide(
        prs,
        "Padrões e outliers",
        [
            "Estados com maior participação do etanol aparecem mais afastados dos mercados dependentes de gasolina.",
            "Outliers combinam grande volume, participação elevada do etanol ou variações mensais fortes.",
            "Os padrões observados fazem sentido para a hipótese de diferenças regionais de substituição.",
        ],
        OUT_FIG / "outliers_pca.png",
    )
    add_slide(
        prs,
        "Comparação e conclusão",
        [
            "PCA explica a contribuição das variáveis e pode apoiar redução de features.",
            "MDS é mais adequado para visualização de similaridade e leitura de vizinhanças.",
            "Os métodos ajudaram a mapear grupos, tendências e registros atípicos úteis para decisões regionais.",
        ],
    )
    prs.save(SLIDES_DIR / "apresentacao_pca_mds_anp.pptx")


def gerar_roteiro_slides_md(ctx: dict) -> None:
    """Cria um apoio simples para a apresentação oral."""
    maior_etanol_txt = ", ".join(
        f"{row.uf} ({row.participacao_media_etanol:.1%})"
        for row in ctx["maior_etanol"].itertuples()
    )
    roteiro = f"""# Roteiro da apresentação: PCA e MDS em dados da ANP

---

## Slide 1: Abertura

**Conteúdo visível no slide:**
- Título: PCA e MDS em dados da ANP
- Subtítulo: Gasolina C, etanol hidratado, preços e volumes por UF

Começar situando o tema. O trabalho usa dados públicos da ANP para observar o preço da gasolina C, o volume vendido de gasolina C e a participação do etanol hidratado nos estados brasileiros ao longo do período {ctx["periodo"]}.

---

## Slide 2: Problema investigado

**Conteúdo visível no slide:**
- Entender associações entre preço da gasolina C, volume vendido e participação do etanol.
- Apoiar decisões de estoque, mix comercial, campanhas e metas regionais.
- Tratar a análise como exploração visual e estatística, sem inferir causalidade direta.

Explicar o cenário da rede de postos ou distribuidora. A gestão precisa tomar decisões de estoque, mix de produtos, campanhas e metas regionais, mas olhar apenas o faturamento agregado não mostra se a mudança de volume veio de preço, substituição por etanol ou diferença regional.

Fala sugerida: "Nossa pergunta principal foi entender em quais estados e períodos o preço da gasolina C se relaciona com o volume vendido e com a presença do etanol."

---

## Slide 3: Dataset e preparação

**Conteúdo visível no slide:**
- Recorte: {ctx["periodo"]}, {ctx["ufs"]} UFs e {ctx["n_registros"]} registros UF-mês.
- Fontes: série histórica mensal de preços e vendas mensais de derivados e etanol da ANP.
- Features padronizadas com StandardScaler para equilibrar escalas diferentes.

Apresentar as duas bases:

- Série histórica mensal de preços por estado.
- Vendas mensais de derivados de petróleo e etanol por UF.

O join foi feito por mês e sigla `uf`. Regiões como Centro-Oeste deixaram de causar falha no merge por grafia diferente entre as bases. Cada linha é um mês em uma UF. Features numéricas receberam `StandardScaler`.

---

## Slide 4: PCA

**Conteúdo visível no slide:**
- PC1 explicou {ctx["variancia_pc1"]:.1%} e PC2 explicou {ctx["variancia_pc2"]:.1%}.
- Variância acumulada dos dois primeiros componentes: {ctx["variancia_total"]:.1%}.
- Principais influências em PC1: {", ".join(ctx["pc1_top"])}.
- Principais influências em PC2: {", ".join(ctx["pc2_top"])}.
- **Imagem:** gráfico de dispersão PCA 2D com pontos coloridos por região (outputs/figures/pca_2d_regiao.png).

Explicar que o PCA transforma as variáveis originais em componentes principais. Os dois primeiros componentes explicaram {ctx["variancia_total"]:.1%} da variância total.

Fala sugerida: "O primeiro eixo ficou ligado principalmente à composição do consumo e ao peso do etanol. O segundo eixo capturou mais as oscilações mensais de preço e volume."

---

## Slide 5: MDS

**Conteúdo visível no slide:**
- Representa proximidades entre registros usando distâncias no espaço padronizado.
- Stress do MDS: {ctx["stress_mds"]:.2f}.
- Registros próximos indicam comportamento semelhante de preço, volume e participação do etanol.
- **Imagem:** gráfico de dispersão MDS 2D com pontos coloridos por região (outputs/figures/mds_2d_regiao.png).

Explicar que o MDS posiciona registros semelhantes próximos entre si, usando distâncias no espaço padronizado. O stress encontrado foi {ctx["stress_mds"]:.2f}; o gráfico deve ser lido como apoio visual, não como prova exata de causalidade.

Fala sugerida: "No MDS, o interesse é ver quem ficou perto de quem. Se dois registros aparecem próximos, eles têm perfis parecidos considerando preço, volume, participação do etanol e variações mensais."

---

## Slide 6: Padrões e outliers

**Conteúdo visível no slide:**
- Estados com maior participação do etanol aparecem mais afastados dos mercados dependentes de gasolina.
- Outliers combinam grande volume, participação elevada do etanol ou variações mensais fortes.
- Os padrões observados fazem sentido para a hipótese de diferenças regionais de substituição.
- **Imagem:** gráfico de barras com os registros mais afastados na projeção PCA (outputs/figures/outliers_pca.png).

Comentar que os estados com maior participação média do etanol no recorte foram {maior_etanol_txt}. São Paulo apareceu com destaque entre os outliers porque combina escala muito alta de vendas e participação forte do etanol.

Fala sugerida: "O ponto atípico não significa erro. Nesse caso, ele mostra que São Paulo tem um comportamento muito diferente em escala e composição do consumo."

---

## Slide 7: Comparação e conclusão

**Conteúdo visível no slide:**
- PCA explica a contribuição das variáveis e pode apoiar redução de features.
- MDS é mais adequado para visualização de similaridade e leitura de vizinhanças.
- Os métodos ajudaram a mapear grupos, tendências e registros atípicos úteis para decisões regionais.

Fechar comparando os métodos:

- PCA ajuda a explicar quais variáveis organizam os dados.
- MDS ajuda a visualizar registros parecidos.
- PCA pode apoiar redução de features em modelos futuros.
- MDS é melhor como visualização exploratória.

Conclusão para falar: "Os dois métodos ajudaram a enxergar que o mercado não se resume ao preço da gasolina. Volume, participação do etanol, preço relativo e variações mensais também estruturam os padrões regionais."
"""
    (SLIDES_DIR / "roteiro_apresentacao_pca_mds_anp.md").write_text(roteiro, encoding="utf-8")


def gerar_notebook() -> None:
    """O arquivo `notebooks/analise_pca_mds_anp.ipynb` é mantido no repositório e não é sobrescrito por esta pipeline.

    Execute o notebook no Jupyter quando quiser repetir os gráficos; o período de análise está na chamada a ``preparar_dados(...)``.
    """


if __name__ == "__main__":
    main()
